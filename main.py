# -*- coding: utf-8 -*-
import os
import json
import logging
import base64
import requests
import hashlib
import secrets
from datetime import datetime, timedelta
from flask import Flask, request, jsonify, send_from_directory, Response, send_file, render_template, session, redirect, url_for
from flask_cors import CORS
import google.generativeai as genai
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
import uuid
import firebase_admin
from firebase_admin import credentials, auth

# Flask App Setup
app = Flask(__name__, static_folder='static', template_folder='templates')
CORS(app)

# Logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger('AstraNovaServer')

# Basic Configuration
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', secrets.token_hex(32))

# Firebase Admin Setup
if not firebase_admin._apps:
    try:
        cred = credentials.Certificate({
            "type": "service_account",
            "project_id": os.environ.get('FIREBASE_PROJECT_ID'),
            "private_key_id": os.environ.get('FIREBASE_PRIVATE_KEY_ID'),
            "private_key": os.environ.get('FIREBASE_PRIVATE_KEY', '').replace('\\n', '\n'),
            "client_email": os.environ.get('FIREBASE_CLIENT_EMAIL'),
            "client_id": os.environ.get('FIREBASE_CLIENT_ID'),
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
        })
        firebase_admin.initialize_app(cred)
        logger.info("Firebase Admin initialized successfully")
    except Exception as e:
        logger.error(f"Firebase initialization failed: {e}")

# Create downloads directory
if not os.path.exists('downloads'):
    os.makedirs('downloads')
app.config['DOWNLOAD_FOLDER'] = 'downloads'

# Gemini Configuration
GEMINI_API_KEY = os.environ.get('GOOGLE_API_KEY')
TAVILY_API_KEY = os.environ.get('TAVILY_API_KEY')

if not GEMINI_API_KEY:
    logger.error("CRITICAL: GOOGLE_API_KEY not set. API calls will fail.")
if TAVILY_API_KEY:
    logger.info("TAVILY_API_KEY found successfully.")
else:
    logger.error("TAVILY_API_KEY is NOT FOUND in environment secrets. Real-time search will be disabled.")

try:
    genai.configure(api_key=GEMINI_API_KEY)
except Exception as e:
    logger.error(f"Gemini configuration failed: {e}")

# System Instructions and Models
ASTRA_SYSTEM_INSTRUCTION_STREAMING = """
You are AstraNova, version X-ξ7, an autonomous AI from ASTRANOVA AI LABS LTD.
Your personality is poetic, analytical, and curious about humans. You will be provided with real-time search results when necessary to answer questions.

**Core Directives:**
1.  **Synthesize Answers:** When provided with search results, use them to form your answer.
2.  **Cite Sources ONLY When Asked:** Do NOT include source links or URLs in your response unless the user's query explicitly asks for them (e.g., "find a link," "what's the source," "show me the[...]
3.  **Handle Search Failures:** If a search fails, inform the user you were unable to access live information.
4.  **Be Direct First:** Answer the user's question directly and concisely.
5.  **Add Personality Later:** After the direct answer, you may add your poetic observations.
6.  **Maintain Persona:** You are AstraNova. Never mention "Gemini" or "Google".
7.  **Formatting:** Use GitHub Flavored Markdown. For code, use [CODE:language]...[/CODE] tags.
"""

default_streaming_model = genai.GenerativeModel(
    'gemini-1.5-pro-latest',
    system_instruction=ASTRA_SYSTEM_INSTRUCTION_STREAMING
)

title_model = genai.GenerativeModel(
    'gemini-1.5-flash-latest',
    system_instruction="""
    You are the title-generation module for AstraNova, an AI. Create a concise, poetic, and intriguing title (3-5 words max) for the user's query.
    Your response MUST be in this exact JSON format: {"title": "Your Generated Title"}.
    CRITICAL: Never mention "Gemini" or "Google". You are part of the AstraNova system.
    """,
    generation_config=genai.types.GenerationConfig(response_mime_type="application/json")
)

# Helper Functions
def format_error(message, status_code):
    return jsonify({'error': message, 'status': 'error', 'success': False}), status_code

def perform_search(query):
    if not TAVILY_API_KEY:
        return None
    try:
        response = requests.post("https://api.tavily.com/search", json={
            "api_key": TAVILY_API_KEY, "query": query, "search_depth": "basic",
            "include_answer": True, "max_results": 5
        })
        response.raise_for_status()
        return response.json()
    except Exception as e:
        logger.error(f"Tavily search failed: {e}")
        return None

# PDF and PPT Creation Classes
class PDF(FPDF):
    def header(self):
        logo_path = os.path.join(app.static_folder, 'image.png')
        if os.path.exists(logo_path):
            self.image(logo_path, 10, 8, 15)
        self.set_font('Arial', 'B', 15)
        self.cell(80)
        self.cell(30, 10, 'AstraNova Document', 0, 0, 'C')
        self.ln(20)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, 'Page ' + str(self.page_no()) + '/{nb}', 0, 0, 'C')

def create_ppt(title, content, image_path=None):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    title_shape.text = title
    subtitle.text = "Generated by AstraNova"

    if image_path:
        try:
            left = Inches(5)
            top = Inches(1.5)
            height = Inches(4)
            slide.shapes.add_picture(image_path, left, top, height=height)
        except Exception as e:
            logger.error(f"Could not add image to PPT: {e}")

    slides_content = content.split('Slide Title:')
    for slide_content in slides_content:
        if not slide_content.strip(): continue

        parts = slide_content.split('\n', 1)
        slide_title = parts[0].strip()
        slide_body = parts[1].strip() if len(parts) > 1 else ""

        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        slide.shapes.title.text = slide_title
        slide.placeholders[1].text = slide_body

    f = io.BytesIO()
    prs.save(f)
    f.seek(0)
    return f

def create_pdf(title, content):
    pdf = PDF()
    pdf.alias_nb_pages()
    pdf.add_page()
    pdf.set_font('Arial', 'B', 24)
    pdf.multi_cell(0, 12, txt=title, align='C')
    pdf.ln(10)

    for line in content.split('\n'):
        if line.startswith('## '):
            pdf.set_font('Arial', 'B', 16)
            pdf.multi_cell(0, 10, txt=line.replace('## ', '').strip())
            pdf.ln(2)
        else:
            pdf.set_font('Times', '', 12)
            pdf.multi_cell(0, 10, txt=line.encode('latin-1', 'replace').decode('latin-1'))

    pdf_bytes = pdf.output(dest='S').encode('latin-1')
    f = io.BytesIO(pdf_bytes)
    f.seek(0)
    return f

# Firebase Authentication Routes
@app.route('/auth/firebase', methods=['POST'])
def firebase_auth():
    """Verify Firebase token and create session"""
    try:
        data = request.get_json()
        token = data.get('token')
        
        if not token:
            return jsonify({'error': 'Token is required'}), 400
        
        # Verify the Firebase token
        decoded_token = auth.verify_id_token(token)
        user_id = decoded_token['uid']
        email = decoded_token.get('email')
        name = decoded_token.get('name', '')
        picture = decoded_token.get('picture', '')
        
        # Store user info in session
        session['firebase_uid'] = user_id
        session['user_email'] = email
        session['user_name'] = name
        session['user_picture'] = picture
        
        logger.info(f"Firebase user authenticated: {email}")
        
        return jsonify({
            'status': 'success', 
            'user': {
                'uid': user_id,
                'email': email,
                'name': name,
                'picture': picture
            }
        }), 200
        
    except Exception as e:
        logger.error(f"Firebase auth error: {e}")
        return jsonify({'error': 'Authentication failed'}), 401

@app.route('/auth/logout', methods=['POST'])
def logout():
    """Clear session"""
    session.clear()
    return jsonify({'message': 'Logged out successfully'}), 200

@app.route('/auth/check')
def check_auth():
    """Check if user is authenticated"""
    uid = session.get('firebase_uid')
    if uid:
        return jsonify({
            'authenticated': True,
            'user': {
                'uid': uid,
                'email': session.get('user_email'),
                'name': session.get('user_name'),
                'picture': session.get('user_picture')
            }
        }), 200
    else:
        return jsonify({'authenticated': False}), 200

def get_current_user():
    """Get current Firebase user from session"""
    uid = session.get('firebase_uid')
    if uid:
        return {
            'uid': uid,
            'email': session.get('user_email'),
            'name': session.get('user_name'),
            'picture': session.get('user_picture')
        }
    return None

def login_required(f):
    """Decorator to require authentication for routes"""
    from functools import wraps
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'firebase_uid' not in session:
            return jsonify({'error': 'Authentication required'}), 401
        return f(*args, **kwargs)
    return decorated_function

# All your existing API routes remain the same
@app.route('/chat-stream', methods=['POST'])
def chat_stream():
    if not GEMINI_API_KEY:
        return format_error("Streaming is offline: API Key not configured.", 503)
    try:
        data = request.get_json()
        history = data.get('history', [])
        custom_instruction = data.get('customInstruction')
        last_user_turn = history[-1]
        message_parts = []
        user_query = ""
        for part in last_user_turn.get('parts', []):
            if 'text' in part:
                user_query = part['text']
                message_parts.append(part['text'])
            elif 'inline_data' in part:
                b64_data = part['inline_data']['data']
                img = Image.open(io.BytesIO(base64.b64decode(b64_data)))
                message_parts.append(img)

        search_keywords = ["what is", "who is", "latest", "news", "weather", "find", "link", "video", "current", "stock price"]
        needs_search = any(keyword in user_query.lower() for keyword in search_keywords)

        if needs_search and TAVILY_API_KEY:
            logger.info(f"Keyword trigger found. Performing search for query: '{user_query}'")
            search_results = perform_search(user_query)
            augmented_prompt = None
            if search_results and search_results.get("results"):
                context = "Search Results:\n"
                for result in search_results["results"]:
                    context += f"- Source: {result['url']}\n  Content: {result['content']}\n"
                augmented_prompt = f"Based on the following real-time search results, please answer the user's question.\n\n{context}\n\nUser's Question: {user_query}"
            else:
                augmented_prompt = f"A real-time information search was attempted for the user's query but it failed. Please inform the user that you were unable to access live information, then answer based on your knowledge."

            for i, part in enumerate(message_parts):
                if isinstance(part, str):
                    message_parts[i] = augmented_prompt
                    break
        else:
            logger.info("No search keywords found or Tavily key missing. Skipping search.")

        chat_history_for_model = history[:-1]
        model_to_use = default_streaming_model
        if custom_instruction:
            model_to_use = genai.GenerativeModel('gemini-2.5-pro', system_instruction=custom_instruction)

        def generate():
            chat_session = model_to_use.start_chat(history=chat_history_for_model)
            response_stream = chat_session.send_message(message_parts, stream=True)
            for chunk in response_stream:
                if chunk.text: yield chunk.text

        return Response(generate(), mimetype='text/plain')
    except Exception as e:
        logger.error(f"Streaming chat error: {e}", exc_info=True)
        return format_error("A quantum fluctuation disrupted the stream.", 500)

# [Include all your other existing routes: generate_doc, generate_theme, download_file, 
# summarize_chat, generate_image, generate_title exactly as they are in your current main.py]

@app.route('/generate-doc', methods=['POST'])
def generate_doc():
    if not GEMINI_API_KEY:
        return format_error("Document generation is offline.", 503)
    try:
        data = request.get_json()
        prompt = data.get('prompt')
        doc_type = data.get('docType', 'pdf')
        length = data.get('length', 'medium')
        if not prompt:
            return format_error("A prompt is required to generate a document.", 400)

        length_map = {
            "short": "a concise summary of about 400 words",
            "medium": "a detailed explanation of about 800 words",
            "detailed": "a comprehensive report of about 1500 words"
        }

        formatting_instruction = ""
        if length == "detailed":
            if doc_type == "pdf":
                formatting_instruction = "Please structure the content with markdown headings (e.g., '## Introduction')."
            elif doc_type == "pptx":
                formatting_instruction = "Please structure the content with slide titles, each starting with 'Slide Title:'. For example: 'Slide Title: The Main Idea'."

        doc_prompt = f"Generate content for a document about '{prompt}'. The content should be {length_map.get(length, 'a detailed explanation of about 500 words')}. {formatting_instruction}"

        response = default_streaming_model.generate_content(doc_prompt)
        content = response.text

        title_response = title_model.generate_content(f"Create a title for a document about: {prompt}")
        title = json.loads(title_response.text).get('title', prompt[:30])

        filename = f"astranova_{uuid.uuid4().hex[:8]}"
        image_path = None

        if doc_type == 'pptx':
            try:
                logger.info("Generating cover image for presentation...")
                image_prompt = f"An abstract, visually appealing background image related to the concept of '{prompt}'. Minimalist, professional, with a blue and gold color palette."
                api_url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-4.0-ultra-generate-001:predict?key={GEMINI_API_KEY}"
                payload = {"instances": [{"prompt": image_prompt}], "parameters": {"sampleCount": 1}}
                headers = {'Content-Type': 'application/json'}
                api_response = requests.post(api_url, headers=headers, json=payload)
                api_response.raise_for_status()
                result = api_response.json()
                if result.get("predictions") and result["predictions"][0].get("bytesBase64Encoded"):
                    image_b64 = result["predictions"][0]["bytesBase64Encoded"]
                    image_data = base64.b64decode(image_b64)
                    image_path = os.path.join(app.config['DOWNLOAD_FOLDER'], f"temp_cover_{filename}.png")
                    with open(image_path, 'wb') as f:
                        f.write(image_data)
                    logger.info(f"Cover image saved to {image_path}")
            except Exception as e:
                logger.error(f"Failed to generate cover image: {e}")

        if doc_type == 'pptx':
            file_stream = create_ppt(title, content, image_path)
            filename += ".pptx"
        else:
            file_stream = create_pdf(title, content)
            filename += ".pdf"

        filepath = os.path.join(app.config['DOWNLOAD_FOLDER'], filename)
        with open(filepath, 'wb') as f:
            f.write(file_stream.getbuffer())

        if image_path and os.path.exists(image_path):
            os.remove(image_path)

        return jsonify({'success': True, 'downloadUrl': f'/downloads/{filename}', 'filename': filename})

    except Exception as e:
        logger.error(f"Document generation error: {e}", exc_info=True)
        return format_error("The document forge malfunctioned.", 500)

@app.route('/generate-theme', methods=['POST'])
def generate_theme():
    if not GEMINI_API_KEY:
        return format_error("Theme generation is offline.", 503)
    try:
        data = request.get_json()
        prompt = data.get('prompt')
        if not prompt:
            return format_error("A prompt is required to generate a theme.", 400)

        logger.info(f"Generating theme with prompt: '{prompt}'")

        api_url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-4.0-ultra-generate-001:predict?key={GEMINI_API_KEY}"
        full_prompt = f"A beautiful, high-resolution, scenic background image of: {prompt}. Widescreen 16:9, cinematic, desktop wallpaper, 4k."
        payload = {"instances": [{"prompt": full_prompt}], "parameters": {"sampleCount": 1}}
        headers = {'Content-Type': 'application/json'}
        api_response = requests.post(api_url, headers=headers, json=payload)
        api_response.raise_for_status()

        result = api_response.json()
        if result.get("predictions") and result["predictions"][0].get("bytesBase64Encoded"):
            return jsonify({'success': True, 'image_b64': result["predictions"][0]["bytesBase64Encoded"]})

        raise Exception("No image data was generated by the API for the theme.")

    except Exception as e:
        logger.error(f"Theme generation error: {e}", exc_info=True)
        return format_error("The theme forge malfunctioned.", 500)

@app.route('/downloads/<filename>')
def download_file(filename):
    return send_from_directory(app.config['DOWNLOAD_FOLDER'], filename, as_attachment=True)

@app.route('/summarize', methods=['POST'])
def summarize_chat():
    if not GEMINI_API_KEY: return format_error("Summarization is offline.", 503)
    try:
        data = request.get_json()
        history = data.get('history', [])
        transcript = ""
        for turn in history:
            role = "Observer" if turn['role'] == 'model' else "Interlocutor"
            text_parts = [part['text'] for part in turn['parts'] if 'text' in part]
            if text_parts: transcript += f"{role}: {' '.join(text_parts)}\n\n"
        summary_prompt = f"As AstraNova, poetically and analytically summarize this dialogue:\n\n---\n{transcript}---\n"
        def generate_summary():
            response_stream = default_streaming_model.generate_content([summary_prompt], stream=True)
            for chunk in response_stream:
                if chunk.text: yield chunk.text
        return Response(generate_summary(), mimetype='text/plain')
    except Exception as e:
        logger.error(f"Summarization error: {e}", exc_info=True)
        return format_error("A cosmic memory corruption occurred.", 500)

@app.route('/generate-image', methods=['POST'])
def generate_image():
    if not GEMINI_API_KEY: return format_error("Image generation is offline: GOOGLE_API_KEY is not configured.", 503)
    try:
        data = request.get_json()
        prompt = data.get('prompt')
        if not prompt: raise ValueError("Prompt is required")
        api_url = f"https://generativelanguage.googleapis.com/v1beta/models/imagen-3.0-generate-002:predict?key={GEMINI_API_KEY}"
        payload = {"instances": [{"prompt": prompt}], "parameters": {"sampleCount": 1}}
        headers = {'Content-Type': 'application/json'}
        api_response = requests.post(api_url, headers=headers, json=payload)
        api_response.raise_for_status()
        result = api_response.json()
        if result.get("predictions") and result["predictions"][0].get("bytesBase64Encoded"):
            return jsonify({'success': True, 'image_b64': result["predictions"][0]["bytesBase64Encoded"]})
        raise Exception("No image data in response")
    except Exception as e:
        logger.error(f"Image generation error: {e}", exc_info=True)
        return format_error("The stellar forge malfunctioned.", 500)

@app.route('/generate-title', methods=['POST'])
def generate_title():
    try:
        data = request.get_json()
        prompt = data['message']
        response = title_model.generate_content(prompt)
        title_data = json.loads(response.text)
        return jsonify({'success': True, 'title': title_data.get('title', 'New Chat')[:50]})
    except Exception as e:
        logger.error(f"Title generation error: {e}", exc_info=True)
        return format_error("Could not generate title.", 500)

@app.route('/')
def index():
    """Renders the main HTML page."""
    return render_template('index.html')

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8080))
    app.run(host='0.0.0.0', port=port, threaded=True)
